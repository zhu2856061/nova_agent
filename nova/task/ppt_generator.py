# -*- coding: utf-8 -*-
# @Time   : 2025/08/12 10:24
# @Author : zip
# @Moto   : Knowledge comes from decomposition
# https://github.com/HansonJames/langchain_universal_tools
import json
import os
import sys
from io import BytesIO
from typing import Dict, List, Optional

import requests
from dotenv import load_dotenv
from langchain.agents import AgentExecutor, create_openai_functions_agent
from langchain.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain.tools import Tool
from langchain_openai import ChatOpenAI
from openai import OpenAI
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pydantic import BaseModel, Field
from tavily import TavilyClient

# 加载环境变量
load_dotenv()


# 配置日志输出
def log_step(message: str):
    """打印生成步骤"""
    print(f"\033[94m[生成进度]\033[0m {message}", file=sys.stderr)


class PPTGenerator:
    """PPT生成器"""

    def __init__(self):
        """初始化PPT生成器"""
        # 初始化语言模型
        self.llm = ChatOpenAI(model_name="gpt-4o", temperature=0)

        # 初始化OpenAI客户端
        self.client = OpenAI()

        # 初始化工具
        tavily_client = TavilyClient()

        def search_with_tavily(query: str) -> str:
            """使用Tavily搜索信息"""
            try:
                search_result = tavily_client.search(query, search_depth="advanced")
                # 提取前5个结果的内容
                # 限制每个结果的内容长度
                contents = []
                for result in search_result.get("results", [])[:3]:  # 只取前3个结果
                    content = result.get("content", "")
                    if len(content) > 500:  # 限制每个内容最多500字符
                        content = content[:500] + "..."
                    contents.append(content)
                return " ".join(contents)
            except Exception as e:
                print(f"\033[91m[搜索错误]\033[0m {str(e)}", file=sys.stderr)
                return ""

        self.tools = [
            Tool(
                name="search",
                func=search_with_tavily,
                description="用于搜索最新的市场数据、行业趋势和相关信息",
            )
        ]

        # 创建提示模板
        self.prompt = ChatPromptTemplate.from_messages(
            [
                (
                    "system",
                    """你是一个专业的PPT设计师，可以帮助用户生成高质量的PPT。你必须严格按照以下步骤操作：

1. 从用户输入中提取两个关键信息：
   - 主题：通常在引号中，如"2025年的创业机遇"
   - 页数：通常以"需要X页"的形式出现，如"需要3页"表示pages=3

2. 调用generate_ppt工具时，必须同时提供topic和pages两个参数。

3. 重要提醒：
   - 你必须从用户输入中找到"需要X页"，并将X作为pages参数的值
   - 你必须同时提供topic和pages两个参数
   - 你绝对不能省略任何参数
   - 你绝对不能使用默认的页数

4. 在调用工具之前，请先确认：
   - 你是否已经从用户输入中找到了页数？
   - 你是否同时准备了两个参数？
   如果有任何一个问题的答案是"否"，请不要调用工具。""",
                ),
                ("human", "{input}"),
                MessagesPlaceholder(variable_name="agent_scratchpad"),
            ]
        )

        # 创建代理
        self.agent = create_openai_functions_agent(
            llm=self.llm, tools=self.tools, prompt=self.prompt
        )

        self.agent_executor = AgentExecutor(
            agent=self.agent, tools=self.tools, verbose=True, handle_parsing_errors=True
        )

    def _clean_json_string(self, text: str) -> str:
        """清理JSON字符串"""
        # 找到第一个完整的JSON对象
        start = text.find("{")
        end = -1
        brace_count = 0
        for i in range(start, len(text)):
            if text[i] == "{":
                brace_count += 1
            elif text[i] == "}":
                brace_count -= 1
                if brace_count == 0:
                    end = i + 1
                    break

        if start == -1 or end == -1:
            raise ValueError("未找到有效的JSON内容")

        json_str = text[start:end]

        # 移除代码块标记和注释
        json_str = json_str.replace("```json", "")
        json_str = json_str.replace("```", "")
        json_str = json_str.replace("/*", "")
        json_str = json_str.replace("*/", "")

        # 添加缺失的逗号
        lines = json_str.split("\n")
        processed_lines = []
        for i, line in enumerate(lines):
            line = line.strip()
            if not line:
                continue

            # 如果不是最后一行且当前行以引号结尾且下一行不是结束标记
            if (
                i < len(lines) - 1
                and line.rstrip().endswith('"')
                and not lines[i + 1].strip().startswith("}")
                and not lines[i + 1].strip().startswith("]")
            ):
                line += ","

            # 如果不是最后一行且当前行以大括号或方括号结尾且下一行不是结束标记
            if (
                i < len(lines) - 1
                and (line.rstrip().endswith("}") or line.rstrip().endswith("]"))
                and not lines[i + 1].strip().startswith("}")
                and not lines[i + 1].strip().startswith("]")
            ):
                line += ","

            processed_lines.append(line)

        json_str = "\n".join(processed_lines)

        # 确保所有属性名都有双引号
        for prop in ["title", "slides", "content", "image_prompt", "search_query"]:
            json_str = json_str.replace(f"{prop}:", f'"{prop}":')
            json_str = json_str.replace(f'"{prop}" :', f'"{prop}":')

        # 确保所有字符串值都有双引号
        import re

        json_str = re.sub(r':\s*([^",\{\[\]\}\s][^",\{\[\]\}\s]*)', r': "\1"', json_str)

        return json_str

    def generate_outline(self, topic: str, pages: int = 3) -> dict:
        """生成PPT大纲"""
        log_step(f"开始为主题「{topic}」生成{pages}页PPT大纲...")

        # 生成提示
        prompt = f"""请为主题「{topic}」生成一个完整的PPT大纲，要求生成{pages}页。

请严格按照以下JSON格式返回，注意：
1. 所有属性和数组项之间必须用英文逗号分隔
2. 所有字符串必须用英文双引号包裹
3. 不要在JSON中添加任何注释
4. 确保生成的是有效的JSON格式

{{
    "title": "PPT标题",
    "slides": [
        {{
            "title": "第一页标题",
            "content": [
                "要点1",
                "要点2",
                "要点3"
            ],
            "image_prompt": "配图描述",
            "search_query": "搜索关键词"
        }}
    ]
}}

内容要求：
1. 必须生成{pages}页幻灯片，内容要丰富全面
2. 每页幻灯片需包含：
   - 标题：清晰的中文标题（不超过15字）
   - 内容：3-4个简洁的要点（每点不超过12字）
   - 图片描述：用于生成配图的详细描述
   - 搜索关键词：用于获取最新数据
3. 结构要完整，包含：
   - 开篇：引言和背景介绍
   - 主体：详细分析和论述
   - 结尾：总结和展望
4. 内容要专业，数据要准确
5. 表达要简洁，逻辑要清晰"""

        response = self.agent_executor.invoke({"input": prompt})

        log_step("大纲生成完成，开始解析内容...")
        try:
            json_str = self._clean_json_string(response["output"])
            return json.loads(json_str)
        except Exception as e:
            print(f"\033[91m[错误]\033[0m JSON解析失败: {str(e)}", file=sys.stderr)
            print(
                f"\033[93m[调试]\033[0m 原始输出:\n{response['output']}",
                file=sys.stderr,
            )
            raise

    def search_info(self, query: str) -> str:
        """搜索补充信息"""
        log_step(f"正在搜索关键信息：{query}")

        response = self.agent_executor.invoke(
            {
                "input": f"""搜索'{query}'相关的最新数据和信息，生成5个重要观点。要求：

1. 每个观点必须包含具体数字或百分比
2. 每个观点字数在10-50字之间
3. 观点要突出增长趋势和发展前景
4. 内容要涵盖市场规模、增长率、未来预测等方面
5. 按重要性排序，最重要的观点放在前面
6. 每个观点单独成段，使用从1开始的数字序号
7. 语言要简洁专业，数据要准确可信

注意：
- 必须使用从1到5的数字序号，不要使用字母序号
- 每个观点都是独立的，不需要子分类
- 每行只使用一个序号，格式为"1. "，"2. "等

示例格式：
1. 全球市场规模达到1000亿美元，同比增长30%。
2. 中国市场份额占40%，预计2025年突破60%。
3. 相关产业带动就业增长25%，创造150万个新岗位。
4. 技术创新投入增加35%，推动产业升级。
5. 国际合作项目增长20%，促进全球化发展。"""
            }
        )

        return response["output"]

    def generate_image(self, prompt: str) -> Optional[BytesIO]:
        """使用DALL-E生成图片"""
        log_step(f"正在生成配图：{prompt}")
        try:
            # 使用DALL-E 3生成图片
            response = self.client.images.generate(
                model="dall-e-3",
                prompt=prompt,
                size="1792x1024",
                quality="standard",
                n=1,
            )

            # 获取图片URL
            image_url = response.data[0].url

            # 下载图片
            image_response = requests.get(image_url)
            if image_response.status_code == 200:
                return BytesIO(image_response.content)
            else:
                raise Exception(f"图片下载失败：{image_response.status_code}")
        except Exception as e:
            print(f"\033[91m[错误]\033[0m 图片生成失败: {str(e)}", file=sys.stderr)
            return None

    def create_ppt(self, outline: dict, output_path: str):
        """创建PPT文件"""
        log_step("开始创建PPT文件...")

        prs = Presentation()

        # 设置16:9比例
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)

        # 创建封面
        log_step("正在创建封面...")
        cover_slide = prs.slides.add_slide(prs.slide_layouts[6])  # 使用空白布局

        # 添加标题
        title_left = Inches(0.5)
        title_top = Inches(0.5)
        title_width = Inches(15)
        title_height = Inches(1.2)

        title_box = cover_slide.shapes.add_textbox(
            title_left, title_top, title_width, title_height
        )

        title_frame = title_box.text_frame
        title_frame.text = outline["title"]
        title_font = title_frame.paragraphs[0].font
        title_font.size = Pt(44)
        title_font.bold = True
        title_font.name = "微软雅黑"
        title_font.color.rgb = RGBColor(31, 73, 125)  # 深蓝色

        # 生成封面图片
        cover_image = self.generate_image(
            f"创建一张体现'{outline['title']}'主题的商务风格封面图"
        )
        if cover_image:
            left = Inches(1)
            top = Inches(2)
            width = Inches(14)
            height = Inches(6)
            cover_slide.shapes.add_picture(cover_image, left, top, width, height)

        # 创建内容页
        log_step("开始创建内容页...")
        for i, slide_data in enumerate(outline["slides"], 1):
            log_step(f"正在创建第 {i} 页：{slide_data['title']}")

            # 使用空白布局
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # 添加标题
            title_left = Inches(0.5)
            title_top = Inches(0.5)
            title_width = Inches(15)
            title_height = Inches(1.0)

            title_box = slide.shapes.add_textbox(
                title_left, title_top, title_width, title_height
            )

            title_frame = title_box.text_frame
            title_frame.text = slide_data["title"]
            title_font = title_frame.paragraphs[0].font
            title_font.size = Pt(32)
            title_font.bold = True
            title_font.name = "微软雅黑"
            title_font.color.rgb = RGBColor(31, 73, 125)  # 深蓝色

            # 创建左侧文本区域
            content_left = Inches(0.8)
            content_top = Inches(1.8)
            content_width = Inches(6.5)
            content_height = Inches(5.5)

            content_box = slide.shapes.add_textbox(
                content_left, content_top, content_width, content_height
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP

            # 添加主要内容
            for idx, point in enumerate(slide_data["content"], 1):
                p = tf.add_paragraph()
                p.text = f"{idx}. {point}"  # 使用数字序号
                p.font.size = Pt(20)
                p.font.name = "微软雅黑"
                p.font.color.rgb = RGBColor(0, 0, 0)  # 黑色
                p.space_after = Pt(12)

            # 添加搜索结果
            if "search_query" in slide_data and slide_data["search_query"]:
                search_result = self.search_info(slide_data["search_query"])
                if search_result:
                    # 添加数据分析标题
                    p = tf.add_paragraph()
                    p.text = "\n市场数据分析"
                    p.font.bold = True
                    p.font.size = Pt(24)
                    p.font.name = "微软雅黑"
                    p.font.color.rgb = RGBColor(68, 114, 196)  # 深蓝色
                    p.space_after = Pt(12)

                    # 分点展示搜索结果
                    points = search_result.split("\n")

                    # 根据内容长度设置不同的字体大小
                    point_number = 1  # 初始化序号
                    for point in points:
                        if point.strip():
                            p = tf.add_paragraph()
                            p.text = f"{point.strip()}"  # 使用数字序号
                            p.font.size = Pt(16)  # 统一字体大小
                            p.font.name = "微软雅黑"
                            p.font.color.rgb = RGBColor(0, 0, 0)  # 黑色
                            p.space_after = Pt(10)
                            point_number += 1  # 递增序号

            # 添加配图
            if "image_prompt" in slide_data and slide_data["image_prompt"]:
                image_data = self.generate_image(slide_data["image_prompt"])
                if image_data:
                    image_left = Inches(8)
                    image_top = Inches(1.8)
                    image_width = Inches(7.2)
                    image_height = Inches(5.2)
                    slide.shapes.add_picture(
                        image_data, image_left, image_top, image_width, image_height
                    )

        # 保存PPT
        log_step("正在保存PPT文件...")
        prs.save(output_path)
        log_step(f"PPT生成完成！文件保存在：{output_path}")

    def generate(self, topic: str, pages: int = 3) -> str:
        """生成完整的PPT"""
        # 创建输出目录
        output_dir = "output/ppt"
        os.makedirs(output_dir, exist_ok=True)

        # 生成大纲
        outline = self.generate_outline(topic, pages)

        # 生成文件路径
        output_path = os.path.join(output_dir, f"{topic}.pptx")

        # 创建PPT
        self.create_ppt(outline, output_path)

        return output_path

    def generate_with_logging(self, topic: str, pages: int = None) -> str:
        """生成PPT并记录过程"""
        if pages is None:
            pages = 3  # 默认3页
        print(f"\n\033[94m[工具调用]\033[0m 开始生成{pages}页PPT，主题：{topic}")
        try:
            result = self.generate(topic, pages)
            print(f"\033[92m[工具完成]\033[0m PPT生成成功：{result}")
            return result
        except Exception as e:
            print(f"\033[91m[工具错误]\033[0m PPT生成失败：{str(e)}")
            raise


if __name__ == "__main__":
    # 使用示例
    generator = PPTGenerator()

    # 测试输入
    input_data = """请为我生成一个主题为"2025年的创业机遇"的PPT，需要3页。要求：
    1. 内容要全面且具有前瞻性
    2. 包含最新的市场数据和趋势
    3. 重点分析未来的创业方向
    4. 配图要专业且美观
    """

    print(f"\n\033[94m[用户输入]\033[0m\n{input_data}")

    # 执行测试
    result = generator.agent_executor.invoke({"input": input_data})

    # 提取文件路径
    output = result["output"]
    path_start = output.find("output/ppt")
    path_end = output.find(".pptx") + 5
    if path_start != -1 and path_end != -1:
        file_path = output[path_start:path_end]
        print(f"\n\033[92m[执行结果]\033[0m")
        print(f"PPT已生成并保存到：{file_path}")
    else:
        print(f"\n\033[92m[执行结果]\033[0m")
        print(output)
