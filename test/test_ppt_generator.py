from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# 创建演示文稿对象
prs = Presentation()

# 第一页：标题页
slide_layout = prs.slide_layouts[0]  # 0是标题页布局
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Python-pptx 演示"
subtitle.text = "使用Python创建PowerPoint演示文稿\n日期：2025年9月"

# 第二页：内容页
slide_layout = prs.slide_layouts[1]  # 1是标题+内容布局
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "主要功能"
tf = content.text_frame
tf.text = "python-pptx库的主要功能："

# 添加项目符号
p = tf.add_paragraph()
p.text = "创建新的演示文稿和幻灯片"
p.level = 1  # 项目符号层级

p = tf.add_paragraph()
p.text = "修改现有演示文稿内容"
p.level = 1

p = tf.add_paragraph()
p.text = "设置文本格式、颜色和字体"
p.level = 1

# 第三页：自定义布局页
slide_layout = prs.slide_layouts[5]  # 5是仅标题布局
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "自定义文本框示例"

# 添加自定义文本框
left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(4)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

tf.text = "这是一个自定义文本框"

p = tf.add_paragraph()
p.text = "可以设置不同的字体大小和颜色"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(0x2E, 0x75, 0xB5)  # 蓝色

# 第四页：图片页
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "图片示例"

# 添加图片（请将路径替换为实际图片路径）
# img_path = "example.jpg"  # 替换为你的图片路径
# left = Inches(1)
# top = Inches(1.5)
# height = Inches(5)
# pic = slide.shapes.add_picture(img_path, left, top, height=height)

# 保存演示文稿
prs.save("python_ppt_demo.pptx")
print("PPT创建成功！文件名为：python_ppt_demo.pptx")
